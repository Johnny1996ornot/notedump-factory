from flask import Flask, request, send_file
from pptx import Presentation
import html 
import base64
import time
from io import BytesIO
import fitz  # PyMuPDF

# Import your existing template logic
import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from template import get_template

app = Flask(__name__)

@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return "No file uploaded", 400
    
    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400

    try:
        file_name = file.filename.lower()
        nav, slides = "", ""
        total_pages = 0
        base_w, base_h = 816, 1054
        unique_storage_id = f"{file_name}_{int(time.time())}"

        if file_name.endswith(('.pptx', '.ppt')):
            ppt = Presentation(file)
            total_pages = len(ppt.slides)
            slide_width_emu = ppt.slide_width or 9144000
            scale_w = base_w / slide_width_emu

            def parse_shapes(shapes):
                html_content = ""
                for shape in shapes:
                    try:
                        if shape.shape_type == 6: 
                            html_content += parse_shapes(shape.shapes)
                            continue

                        top = (shape.top * scale_w) if shape.top else 0
                        left = (shape.left * scale_w) if shape.left else 0
                        width = (shape.width * scale_w) if shape.width else 200
                        height = (shape.height * scale_w) if shape.height else 50

                        if shape.shape_type == 13: 
                            img_stream = BytesIO(shape.image.blob)
                            base64_img = base64.b64encode(img_stream.getvalue()).decode()
                            html_content += f'''
                            <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; height:{height}px; z-index:10; transform: translate(0px, 0px);">
                                <img src="data:image/png;base64,{base64_img}" style="width:100%; height:100%; object-fit:contain;">
                            </div>'''

                        elif shape.has_table:
                            table_html = "<table style='width:100%; height:100%; border-collapse: collapse; font-size:12px;' border='1'>"
                            for row in shape.table.rows:
                                table_html += "<tr>"
                                for cell in row.cells:
                                    table_html += f"<td style='padding:5px;'>{html.escape(cell.text)}</td>"
                                table_html += "</tr>"
                            table_html += "</table>"
                            html_content += f'''
                            <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; background:rgba(255,255,255,0.9); z-index:15; transform: translate(0px, 0px);">
                                <div class="content-area">{table_html}</div>
                            </div>'''

                        elif shape.has_text_frame and shape.text.strip():
                            html_text = ""
                            for paragraph in shape.text_frame.paragraphs:
                                p_text = ""
                                for run in paragraph.runs:
                                    r_txt = html.escape(run.text)
                                    fs_style = ""
                                    if hasattr(run.font, 'size') and run.font.size:
                                        pt_size = run.font.size.pt
                                        px_size = max(10, int(pt_size * scale_w * 12700 * 1.33))
                                        fs_style = f"font-size:{px_size}px;"

                                    if getattr(run.font, 'bold', False) == True: r_txt = f"<strong>{r_txt}</strong>"
                                    if getattr(run.font, 'italic', False) == True: r_txt = f"<em>{r_txt}</em>"
                                    if getattr(run.font, 'underline', False) == True: r_txt = f"<u>{r_txt}</u>"

                                    if fs_style: p_text += f"<span style='{fs_style}'>{r_txt}</span>"
                                    else: p_text += r_txt

                                if not p_text.strip(): html_text += "<br>"
                                else: html_text += f"<div>{p_text}</div>"

                            html_content += f'''
                            <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; min-height:{height}px; z-index:20; transform: translate(0px, 0px);">
                                <div class="content-area" style="word-wrap: break-word; white-space: pre-wrap; overflow-wrap: break-word;">{html_text}</div>
                            </div>'''
                    except Exception as e:
                        continue 
                return html_content

            for i, slide in enumerate(ppt.slides):
                title_text = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
                nav += f'<div class="nav-link" id="link-{i}" onclick="goTo(\'{i}\')"><i class="fas fa-bars drag-handle"></i> <span class="nav-text">{html.escape(title_text)}</span></div>'
                slides += f'<div id="p-{i}" class="page {"active" if i==0 else ""}" data-page-height="{base_h}" style="height:{base_h}px;"> '
                slides += parse_shapes(slide.shapes)
                slides += '</div>'

        elif file_name.endswith('.pdf'):
            doc = fitz.open(stream=file.read(), filetype="pdf")
            total_pages = len(doc)
            first_page = doc[0]
            p_scale = base_w / first_page.rect.width if first_page.rect.width > 0 else 1

            for i, page in enumerate(doc):
                page_width = page.rect.width
                page_height = page.rect.height
                p_scale = base_w / page_width if page_width > 0 else 1
                scaled_height = max(1054, int(page_height * p_scale))

                nav += f'<div class="nav-link" id="link-{i}" onclick="goTo(\'{i}\')"><i class="fas fa-bars drag-handle"></i> <span class="nav-text">Page {i+1}</span></div>'
                slides += f'<div id="p-{i}" class="page {"active" if i==0 else ""}" data-page-height="{scaled_height}" style="height:{scaled_height}px;"> '

                blocks = page.get_text("dict")["blocks"]
                html_content = ""

                for b in blocks:
                    bbox = b["bbox"]
                    top = bbox[1] * p_scale
                    left = bbox[0] * p_scale
                    width = (bbox[2] - bbox[0]) * p_scale
                    height = (bbox[3] - bbox[1]) * p_scale

                    if b["type"] == 0: 
                        text_html = ""
                        for line in b["lines"]:
                            line_html = ""
                            for span in line["spans"]:
                                txt = html.escape(span["text"])
                                if not txt.strip():
                                    line_html += " "
                                    continue

                                px_size = max(10, int(span["size"] * p_scale))
                                fs_style = f"font-size:{px_size}px;"

                                if span["flags"] & 16: txt = f"<strong>{txt}</strong>"
                                if span["flags"] & 2: txt = f"<em>{txt}</em>"

                                line_html += f"<span style='{fs_style}'>{txt}</span>"

                            if not line_html.strip(): text_html += "<br>"
                            else: text_html += f"<div>{line_html}</div>"

                        html_content += f'''
                        <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; min-height:{height}px; z-index:20; transform: translate(0px, 0px);">
                            <div class="content-area" style="word-wrap: break-word; white-space: pre-wrap; overflow-wrap: break-word; font-family: sans-serif;">{text_html}</div>
                        </div>'''

                    elif b["type"] == 1: 
                        img_bytes = b.get("image")
                        if img_bytes:
                            base64_img = base64.b64encode(img_bytes).decode()
                            html_content += f'''
                            <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; height:{height}px; z-index:10; transform: translate(0px, 0px);">
                                <img src="data:image/png;base64,{base64_img}" style="width:100%; height:100%; object-fit:contain;">
                            </div>'''

                slides += html_content
                slides += '</div>'

        dimension_script = f"""
        <script>
            document.addEventListener("DOMContentLoaded", function() {{
                var cvs = document.getElementById('canvas');
                if(cvs) {{
                    cvs.style.width = '{base_w}px';
                    cvs.setAttribute('data-width', '{base_w}');
                    var wInput = document.getElementById('canvas-w-cm');
                    var hInput = document.getElementById('canvas-h-cm');
                    if(wInput) wInput.value = (Math.round(({base_w} / 37.795) * 10) / 10).toFixed(1);
                    if(hInput) hInput.value = (Math.round((parseInt(cvs.style.height) / 37.795) * 10) / 10).toFixed(1);
                }}
            }});
        </script>
        """

        # Build the final HTML string
        final_html = get_template(total_pages)\
            .replace("{{NAV_LINKS}}", nav)\
            .replace("{{SLIDE_CONTENT}}", dimension_script + slides)\
            .replace("{{VISIBLE_TITLE}}", html.escape(file.filename))\
            .replace("{{STORAGE_ID}}", html.escape(unique_storage_id))

        # Send it back to the user's browser to download
        output = BytesIO()
        output.write(final_html.encode('utf-8'))
        output.seek(0)
        
        return send_file(output, mimetype='text/html', as_attachment=True, download_name=f"NoteDump_{file.filename}.html")

    except Exception as e:
        return f"Error Processing File: {str(e)}", 500

# Required for Vercel Serverless
app = app
