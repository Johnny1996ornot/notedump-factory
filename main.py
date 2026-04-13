import os
import time
from io import BytesIO
import base64
import html
from flask import Flask, request, send_file
from pptx import Presentation
import fitz  # PyMuPDF
from template import get_template

app = Flask(__name__)

# CRITICAL FIX FOR RENDER: Tell Flask to accept uploads up to 250 MB
app.config['MAX_CONTENT_LENGTH'] = 250 * 1024 * 1024 

# --- HTML FRONTEND ---
FRONTEND_HTML = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>NoteDump</title>
    <style>
        body { background-color: #000000; color: white; font-family: 'Segoe UI', sans-serif; display: flex; flex-direction: column; align-items: center; justify-content: center; height: 100vh; margin: 0; position: relative; }
        
        /* 1. Header & Logo Container */
        .hero { text-align: center; margin-bottom: 40px; position: relative; }
        .logo-container { position: relative; display: inline-block; padding: 0 20px; }
        .logo-text { font-size: 55px; font-weight: 800; color: #f8fafc; margin: 0; display: inline-block;}
        
        /* 2. Repositioned Buttons (Now relative to Logo) */
        .header-actions { position: absolute; top: -10px; right: -80px; display: flex; align-items: center; gap: 8px; }
        .guide-btn { color: #94a3b8; text-decoration: none; font-size: 14px; font-weight: bold; border: 1px solid #475569; border-radius: 50%; width: 24px; height: 24px; display: flex; align-items: center; justify-content: center; background: #1e293b; transition: 0.2s; }
        .guide-btn:hover { color: white; border-color: #0ea5e9; background: #0ea5e9;}
        .coffee-btn { color: #0ea5e9; text-decoration: none; font-weight: bold; font-size: 11px; border: 1px solid #0ea5e9; padding: 4px 10px; border-radius: 20px; transition: 0.2s; white-space: nowrap; background: transparent; }
        .coffee-btn:hover { background: rgba(14, 165, 233, 0.15); color: #fff; }

        /* 3. Lowered Tagline & Text */
        .tagline { color: #94a3b8; font-size: 18px; margin-top: 15px; margin-bottom: 5px; }
        .support-text { font-size: 12px; color: #475569; margin-top: 5px; letter-spacing: 1px; }

        /* Modal Settings */
        .modal-window { position: fixed; background-color: rgba(0, 0, 0, 0.85); backdrop-filter: blur(5px); top: 0; right: 0; bottom: 0; left: 0; z-index: 99999; visibility: hidden; opacity: 0; transition: all 0.3s; display: flex; justify-content: center; align-items: center; }
        .modal-window:target { visibility: visible; opacity: 1; }
        .modal-content { background: #0f172a; width: 90%; max-width: 650px; padding: 30px; border-radius: 16px; border: 1px solid #334155; color: #f1f5f9; position: relative; max-height: 80vh; overflow-y: auto; text-align: left;}
        .modal-close { position: absolute; top: 20px; right: 25px; color: #64748b; text-decoration: none; font-size: 28px; font-weight: bold; }
        .modal-content h2 { color: #4f46e5; border-bottom: 1px solid #1e293b; padding-bottom: 10px;}
        .pro-tag { color: #10b981; font-weight: bold; } 

        /* Main UI Box */
        .upload-box { background-color: #0f172a; border: 1px dashed #334155; border-radius: 12px; padding: 40px; text-align: center; width: 450px; box-shadow: 0 10px 30px rgba(0,0,0,0.5); }
        .primary-btn { display: inline-block; background-color: #4f46e5; color: white; border: none; padding: 16px 28px; border-radius: 8px; font-weight: bold; cursor: pointer; font-size: 16px; width: 100%; transition: 0.2s; text-align: center; box-sizing: border-box; }
        .primary-btn:hover { background-color: #4338ca; }
        .primary-btn span { color: #c7d2fe; font-size: 14px; display: block; margin-top: 4px; font-weight: normal;}
        .blank-btn { background-color: #1e293b; border: 1px solid #334155; margin-top: 20px;}
        .blank-btn:hover { background-color: #334155; }
        
        /* Progress & Error */
        .progress-container { width: 100%; background-color: #1e293b; border-radius: 4px; overflow: hidden; margin-bottom: 8px; height: 10px; display: none; }
        .progress-fill { width: 0%; height: 100%; background-color: #10b981; transition: width 0.2s ease-out; }
        .status-text { color: #c7d2fe; font-size: 14px; font-weight: bold; display: none; }
        .error-box { display: none; color: #ef4444; background-color: rgba(239, 68, 68, 0.1); border: 1px solid #ef4444; border-radius: 6px; padding: 10px; margin-top: 20px; font-size: 14px; }
        @keyframes pulse { 0%, 100% { opacity: 1; } 50% { opacity: 0.5; } }
        .processing-pulse { animation: pulse 1.5s infinite; color: #fcd34d !important; }
    </style>
    <script>
        function handleUpload(event) {
            const fileInput = event.target;
            if (!fileInput.files.length) return;

            const file = fileInput.files[0];
            const btnLabel = document.getElementById('upload-label');
            const errorBox = document.getElementById('error-box');
            
            // Save original HTML
            const originalHtml = `📥 convert existing file<br><span>200 mb pdf pdx ppt</span>`;
            
            // Trigger UI Loading State
            btnLabel.style.pointerEvents = 'none';
            btnLabel.style.backgroundColor = '#312e81';
            btnLabel.innerHTML = `
                <div class="progress-container" id="prog-container" style="display: block;">
                    <div class="progress-fill" id="prog-fill"></div>
                </div>
                <div class="status-text" id="stat-text" style="display: block;">Uploading: 0%</div>
            `;
            errorBox.style.display = 'none';

            const formData = new FormData();
            formData.append('file', file);

            const xhr = new XMLHttpRequest();
            xhr.open('POST', '/convert', true);
            xhr.responseType = 'blob';

            // Phase 1: Real-time progress bar tracking
            xhr.upload.onprogress = function(e) {
                if (e.lengthComputable) {
                    const percent = Math.round((e.loaded / e.total) * 100);
                    document.getElementById('prog-fill').style.width = percent + '%';
                    document.getElementById('stat-text').innerText = 'Uploading: ' + percent + '%';
                    
                    if (percent >= 100) {
                        // Switch text once upload finishes and parsing begins
                        document.getElementById('stat-text').innerHTML = 'Processing Document...<br><small style="font-weight: normal;">(Converting slides & images)</small>';
                        document.getElementById('stat-text').classList.add('processing-pulse');
                    }
                }
            };

            // FAIL-SAFE: The instant the request resolves (success or fail), reset the button.
            xhr.onloadend = function() {
                btnLabel.style.backgroundColor = '';
                btnLabel.style.pointerEvents = 'auto';
                btnLabel.innerHTML = originalHtml;
                fileInput.value = '';
            };

            // Phase 2: Handle the server's response
            xhr.onload = function() {
                if (xhr.status === 200) {
                    // Success -> Trigger file download natively in browser
                    const url = window.URL.createObjectURL(xhr.response);
                    const a = document.createElement('a');
                    a.style.display = 'none';
                    a.href = url;
                    a.download = `NoteDump_${file.name}.html`;
                    document.body.appendChild(a);
                    a.click();
                    window.URL.revokeObjectURL(url);
                } else {
                    // Server threw an error (like a 500 or 400)
                    const reader = new FileReader();
                    reader.onload = () => { 
                        errorBox.innerHTML = "Server Error: " + reader.result; 
                        errorBox.style.display = 'block';
                    };
                    reader.readAsText(xhr.response);
                }
            };

            xhr.onerror = () => {
                errorBox.innerHTML = "🚨 Network error or timeout. If the file is massive, the server process may have been killed.";
                errorBox.style.display = 'block';
            };

            xhr.send(formData);
        }

        async function handleBlank(event) {
            event.preventDefault();
            const btn = document.getElementById('blank-btn');
            const original = btn.innerHTML;
            
            btn.style.pointerEvents = 'none';
            btn.style.opacity = '0.7';
            btn.innerHTML = '⏳ Generating...';
            
            try {
                const res = await fetch('/blank');
                if (res.ok) {
                    const blob = await res.blob();
                    const url = window.URL.createObjectURL(blob);
                    const a = document.createElement('a');
                    a.style.display = 'none';
                    a.href = url;
                    a.download = `NoteDump_Blank.html`;
                    document.body.appendChild(a);
                    a.click();
                    window.URL.revokeObjectURL(url);
                }
            } catch(e) {
                console.error(e);
            }
            
            btn.style.pointerEvents = 'auto';
            btn.style.opacity = '1';
            btn.innerHTML = original;
        }
    </script>
</head>
<body>
    
    <div class="hero">
        <div class="logo-container">
            <h1 class="logo-text">📝 NoteDump</h1>
            <div class="header-actions">
                <a href="https://buymeacoffee.com/jpramirez" target="_blank" class="coffee-btn">☕ Support</a>
                <a href="#guide-modal" class="guide-btn">?</a>
            </div>
        </div>
        <p class="tagline">Turning your documents into an interactive notebook</p>
        <p class="support-text">PPTX • PPT • PDF</p>
    </div>

    <div id="guide-modal" class="modal-window">
        <div class="modal-content">
            <a href="#" class="modal-close">&times;</a>
            <h2>📝 Welcome to NoteDump</h2>
            <p>Turning your documents into an interactive notebook.</p>
            <h4>✨ NoteDump Features & Guide</h4>
            <ul>
                <li>Providing an HTML based editor that allows for creating seamless interactive notebooks.</li>
                <li>Add pins to images to correctly identify them.</li>
                <li>Add audio or links that will allow a better review.</li>
                <li>Create reviewers with more freedom and specialized tools.</li>
            </ul>
            <h4>🚀 Why Use NoteDump?</h4>
            <ul>
                <li><span class="pro-tag">100% Offline Capable:</span> Once downloaded, your interactive notebook is a single HTML file that works perfectly without an internet connection.</li>
                <li><span class="pro-tag">Ultimate Privacy:</span> No accounts, no cloud syncing, no subscriptions. Your notes and files stay locally on your device.</li>
                <li><span class="pro-tag">Limitless Annotation:</span> Draw freely, drop sticky notes, create custom tables, and pin interactive markers directly onto your lecture slides.</li>
                <li><span class="pro-tag">Universal Compatibility:</span> Works natively on any modern device (desktop, tablet, or phone) using just a standard web browser.</li>
            </ul>
        </div>
    </div>

    <div class="upload-box">
        <input type="file" id="file" accept=".pdf,.pptx,.ppt" style="display: none;" onchange="handleUpload(event);">
        <label id="upload-label" for="file" class="primary-btn">
            📥 convert existing file<br>
            <span>200 mb pdf pdx ppt</span>
        </label>
        
        <button id="blank-btn" class="primary-btn blank-btn" onclick="handleBlank(event)">📓 Create Blank Notebook</button>
        
        <div id="error-box" class="error-box"></div>
    </div>

</body>
</html>
"""

@app.route("/", methods=["GET"])
def index():
    return FRONTEND_HTML

@app.route("/blank", methods=["GET"])
def create_blank():
    blank_nav = '<div class="nav-link active-nav" id="link-0" onclick="goTo(\'0\')"><i class="fas fa-bars drag-handle"></i> <span class="nav-text">Page 1</span></div>'
    blank_slides = '<div id="p-0" class="page active" data-page-width="816" data-page-height="1054" style="width:816px; height:1054px;"></div>'
    unique_blank_id = f"New_Notebook_{int(time.time())}"
    
    html_out = get_template(1) \
        .replace("{{NAV_LINKS}}", blank_nav) \
        .replace("{{SLIDE_CONTENT}}", blank_slides) \
        .replace("{{VISIBLE_TITLE}}", "New_Notebook") \
        .replace("{{STORAGE_ID}}", unique_blank_id)
        
    return send_file(
        BytesIO(html_out.encode('utf-8')),
        mimetype="text/html",
        as_attachment=True,
        download_name="NoteDump_Blank.html"
    )

@app.route("/convert", methods=["POST"])
def convert_file():
    if 'file' not in request.files:
        return "Error: No file part uploaded", 400
    
    file = request.files['file']
    if file.filename == '':
        return "Error: No selected file", 400
        
    file_name = file.filename.lower()
    file_bytes = file.read()
    
    nav = ""
    slides = ""
    total_pages = 0
    unique_storage_id = f"{file_name}_{int(time.time())}"
    base_w = 816 
    base_h = 1054
    
    try:
        if file_name.endswith(('.pptx', '.ppt')):
            ppt = Presentation(BytesIO(file_bytes))
            total_pages = len(ppt.slides)
            slide_width_emu = ppt.slide_width or 9144000
            scale_w = base_w / slide_width_emu

            def parse_shapes(shapes):
                html_content = ""
                for shape in shapes:
                    try:
                        if shape.shape_type == 6: 
                            html_content += parse_shapes(shape.shapes)
                            continue

                        top = (shape.top * scale_w) if shape.top else 0
                        left = (shape.left * scale_w) if shape.left else 0
                        width = (shape.width * scale_w) if shape.width else 200
                        height = (shape.height * scale_w) if shape.height else 50

                        if shape.shape_type == 13: 
                            img_stream = BytesIO(shape.image.blob)
                            base64_img = base64.b64encode(img_stream.getvalue()).decode()
                            html_content += f'''
                            <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; height:{height}px; z-index:10; transform: translate(0px, 0px);">
                                <img src="data:image/png;base64,{base64_img}" style="width:100%; height:100%; object-fit:contain;">
                            </div>'''

                        elif shape.has_table:
                            table_html = "<table style='width:100%; height:100%; border-collapse: collapse; font-size:12px;' border='1'>"
                            for row in shape.table.rows:
                                table_html += "<tr>"
                                for cell in row.cells:
                                    table_html += f"<td style='padding:5px;'>{html.escape(cell.text)}</td>"
                                table_html += "</tr>"
                            table_html += "</table>"
                            html_content += f'''
                            <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; background:rgba(255,255,255,0.9); z-index:15; transform: translate(0px, 0px);">
                                <div class="content-area">{table_html}</div>
                            </div>'''

                        elif shape.has_text_frame and shape.text.strip():
                            html_text = ""
                            for paragraph in shape.text_frame.paragraphs:
                                p_text = ""
                                for run in paragraph.runs:
                                    r_txt = html.escape(run.text)
                                    fs_style = ""
                                    if hasattr(run.font, 'size') and run.font.size:
                                        pt_size = run.font.size.pt
                                        px_size = max(10, int(pt_size * scale_w * 12700 * 1.33))
                                        fs_style = f"font-size:{px_size}px;"

                                    if getattr(run.font, 'bold', False) == True: r_txt = f"<strong>{r_txt}</strong>"
                                    if getattr(run.font, 'italic', False) == True: r_txt = f"<em>{r_txt}</em>"
                                    if getattr(run.font, 'underline', False) == True: r_txt = f"<u>{r_txt}</u>"

                                    if fs_style: p_text += f"<span style='{fs_style}'>{r_txt}</span>"
                                    else: p_text += r_txt

                                if not p_text.strip(): html_text += "<br>"
                                else: html_text += f"<div>{p_text}</div>"

                            html_content += f'''
                            <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; min-height:{height}px; z-index:20; transform: translate(0px, 0px);">
                                <div class="content-area" style="word-wrap: break-word; white-space: pre-wrap; overflow-wrap: break-word;">{html_text}</div>
                            </div>'''
                    except Exception as e:
                        continue 
                return html_content

            for i, slide in enumerate(ppt.slides):
                title_text = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
                nav += f'<div class="nav-link" id="link-{i}" onclick="goTo(\'{i}\')"><i class="fas fa-bars drag-handle"></i> <span class="nav-text">{html.escape(title_text)}</span></div>'
                slides += f'<div id="p-{i}" class="page {"active" if i==0 else ""}" data-page-height="{base_h}" style="height:{base_h}px;"> '
                slides += parse_shapes(slide.shapes)
                slides += '</div>'

        elif file_name.endswith('.pdf'):
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            total_pages = len(doc)

            first_page = doc[0]
            p_scale = base_w / first_page.rect.width if first_page.rect.width > 0 else 1

            for i, page in enumerate(doc):
                page_width = page.rect.width
                page_height = page.rect.height
                p_scale = base_w / page_width if page_width > 0 else 1
                scaled_height = max(1054, int(page_height * p_scale))

                nav += f'<div class="nav-link" id="link-{i}" onclick="goTo(\'{i}\')"><i class="fas fa-bars drag-handle"></i> <span class="nav-text">Page {i+1}</span></div>'
                slides += f'<div id="p-{i}" class="page {"active" if i==0 else ""}" data-page-height="{scaled_height}" style="height:{scaled_height}px;"> '

                blocks = page.get_text("dict")["blocks"]
                html_content = ""

                for b in blocks:
                    bbox = b["bbox"]
                    top = bbox[1] * p_scale
                    left = bbox[0] * p_scale
                    width = (bbox[2] - bbox[0]) * p_scale
                    height = (bbox[3] - bbox[1]) * p_scale

                    if b["type"] == 0: 
                        text_html = ""
                        for line in b["lines"]:
                            line_html = ""
                            for span in line["spans"]:
                                txt = html.escape(span["text"])
                                if not txt.strip():
                                    line_html += " "
                                    continue

                                px_size = max(10, int(span["size"] * p_scale))
                                fs_style = f"font-size:{px_size}px;"

                                if span["flags"] & 16: txt = f"<strong>{txt}</strong>"
                                if span["flags"] & 2: txt = f"<em>{txt}</em>"

                                line_html += f"<span style='{fs_style}'>{txt}</span>"

                            if not line_html.strip(): text_html += "<br>"
                            else: text_html += f"<div>{line_html}</div>"

                        html_content += f'''
                        <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; min-height:{height}px; z-index:20; transform: translate(0px, 0px);">
                            <div class="content-area" style="word-wrap: break-word; white-space: pre-wrap; overflow-wrap: break-word; font-family: sans-serif;">{text_html}</div>
                        </div>'''

                    elif b["type"] == 1: 
                        img_bytes = b.get("image")
                        if img_bytes:
                            base64_img = base64.b64encode(img_bytes).decode()
                            html_content += f'''
                            <div class="canvas-box" style="top:{top}px; left:{left}px; width:{width}px; height:{height}px; z-index:10; transform: translate(0px, 0px);">
                                <img src="data:image/png;base64,{base64_img}" style="width:100%; height:100%; object-fit:contain;">
                            </div>'''

                slides += html_content
                slides += '</div>'

        dimension_script = f"""
        <script>
            document.addEventListener("DOMContentLoaded", function() {{
                var cvs = document.getElementById('canvas');
                if(cvs) {{
                    cvs.style.width = '{base_w}px';
                    cvs.setAttribute('data-width', '{base_w}');
                    var wInput = document.getElementById('canvas-w-cm');
                    var hInput = document.getElementById('canvas-h-cm');
                    if(wInput) wInput.value = (Math.round(({base_w} / 37.795) * 10) / 10).toFixed(1);
                    if(hInput) hInput.value = (Math.round((parseInt(cvs.style.height) / 37.795) * 10) / 10).toFixed(1);
                }}
            }});
        </script>
        """

        final_html = get_template(total_pages)\
            .replace("{{NAV_LINKS}}", nav)\
            .replace("{{SLIDE_CONTENT}}", dimension_script + slides)\
            .replace("{{VISIBLE_TITLE}}", html.escape(file.filename))\
            .replace("{{STORAGE_ID}}", html.escape(unique_storage_id))
            
        return send_file(
            BytesIO(final_html.encode('utf-8')),
            mimetype="text/html",
            as_attachment=True,
            download_name=f"NoteDump_{file.filename}.html"
        )

    except Exception as e:
        return f"Error processing file: {html.escape(str(e))}", 500

if __name__ == "__main__":
    app.run(debug=True)
